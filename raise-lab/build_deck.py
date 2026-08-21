#!/usr/bin/env python3
"""Build the RAISE Lab introduction deck (16:9 PowerPoint)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ----------------------------------------------------------------- brand ---
NAVY      = RGBColor(0x1A, 0x1A, 0x2E)   # site --primary
NAVY2     = RGBColor(0x16, 0x21, 0x3E)   # site --secondary
CRIMSON   = RGBColor(0xE9, 0x45, 0x60)   # site --accent
AMBER     = RGBColor(0xF3, 0x9C, 0x6B)   # site --accent-soft
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x2D, 0x2D, 0x2D)   # site --text
MUTED     = RGBColor(0x66, 0x66, 0x66)   # site --text-light
PANEL     = RGBColor(0xF4, 0xF5, 0xF8)
RULE      = RGBColor(0xE0, 0xE0, 0xE0)
PALE      = RGBColor(0xB8, 0xBE, 0xCC)

HEAD = "Georgia"    # ships with Office on Mac + Windows
BODY = "Calibri"    # ships with Office on Mac + Windows

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.85)               # left/right margin
CW = W - 2 * M                 # content width

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers ---
def slide(notes=""):
    s = prs.slides.add_slide(BLANK)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def bg(s, color):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    return r


def rect(s, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, lw=1.0):
    r = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        r.fill.background()
    else:
        r.fill.solid()
        r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        r.line.width = Pt(lw)
    r.shadow.inherit = False
    if r.has_text_frame:
        r.text_frame.text = ""
    return r


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, font=BODY,
         align=PP_ALIGN.LEFT, spacing=1.0, anchor=MSO_ANCHOR.TOP,
         space_after=0, caps=False):
    """runs: str, or list of (text, {overrides}) tuples, or list of paragraphs."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        chunks = para if isinstance(para, list) else [(para, {})]
        for chunk in chunks:
            txt, ov = (chunk, {}) if isinstance(chunk, str) else chunk
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = ov.get("font", font)
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.italic = ov.get("italic", False)
            f.color.rgb = ov.get("color", color)
    return tb


def header(s, kicker, title, dark=False):
    """Standard content-slide header: small crimson kicker + serif title."""
    tc = PALE if dark else MUTED
    hc = WHITE if dark else NAVY
    text(s, M, Inches(0.52), CW, Inches(0.3),
         [[(kicker.upper(), {})]], size=11.5, color=CRIMSON, bold=True,
         font=BODY, spacing=1.0)
    text(s, M, Inches(0.86), CW, Inches(0.62), title,
         size=30, color=hc, bold=False, font=HEAD)
    rect(s, M, Inches(1.56), Inches(0.62), Pt(3), fill=CRIMSON)
    return tc


def footer(s, n, dark=False):
    if n is None:
        return
    text(s, M, H - Inches(0.62), CW - Inches(0.4), Inches(0.24),
         "RAISE Lab  ·  HCMUT, VNU-HCM",
         size=9.5, color=(PALE if dark else MUTED), font=BODY)
    text(s, W - M - Inches(0.6), H - Inches(0.62), Inches(0.6), Inches(0.24),
         str(n), size=9.5, color=(PALE if dark else MUTED), font=BODY,
         align=PP_ALIGN.RIGHT)


def rows(s, items, top, row_h, gap=Inches(0.10), numbered=True,
         title_size=15, desc_size=11.5, width=None):
    """Vertical stack of numbered direction cards."""
    width = width or CW
    y = top
    for i, (t, d) in enumerate(items, 1):
        rect(s, M, y, width, row_h, fill=PANEL)
        rect(s, M, y, Pt(3.2), row_h, fill=CRIMSON if i % 2 else AMBER)
        bx = M + Inches(0.26)
        if numbered:
            text(s, bx, y + Inches(0.16), Inches(0.42), Inches(0.34),
                 f"{i:02d}", size=15, color=CRIMSON, bold=True, font=HEAD)
            bx += Inches(0.52)
        tw = width - (bx - M) - Inches(0.3)
        text(s, bx, y + Inches(0.15), tw, Inches(0.3), t,
             size=title_size, color=NAVY, bold=True, font=BODY)
        text(s, bx, y + Inches(0.15) + Pt(title_size + 6), tw,
             row_h - Inches(0.5), d, size=desc_size, color=MUTED,
             font=BODY, spacing=1.12)
        y += row_h + gap




# ============================================================ 1 · TITLE ===
s = slide(
    "Opening. Keep it short and warm, and don't read the slide.\n\n"
    "\"RAISE stands for Reasoning in Artificial Intelligence and Software "
    "Engineering. We're a research lab at HCMUT. The short version of what we "
    "do is this: AI systems now produce a lot of software, and we build the "
    "checks that tell you whether to trust any of it.\"\n\n"
    "Set expectations: about fifteen minutes, three research pillars, then "
    "how to work with us."
)
bg(s, NAVY)
rect(s, W - Inches(4.6), 0, Inches(4.6), H, fill=NAVY2)
rect(s, W - Inches(4.6), 0, Pt(3), H, fill=CRIMSON)
rect(s, M, Inches(1.72), Inches(1.05), Pt(4), fill=CRIMSON)
text(s, M, Inches(2.05), Inches(7.6), Inches(1.5), "RAISE",
     size=88, color=WHITE, font=HEAD)
text(s, M, Inches(3.46), Inches(7.4), Inches(0.9),
     "Reasoning in Artificial Intelligence\nand Software Engineering",
     size=20, color=PALE, font=BODY, spacing=1.25)
rect(s, M, Inches(4.62), Inches(7.4), Pt(1), fill=RGBColor(0x3A, 0x3F, 0x55))
text(s, M, Inches(4.92), Inches(7.4), Inches(0.9),
     [[("Ho Chi Minh City University of Technology (HCMUT), VNU-HCM", {})]],
     size=14, color=WHITE, font=BODY)
text(s, M, Inches(5.30), Inches(7.4), Inches(0.9),
     [[("Dr. Xuan-Bach Le", {"bold": True, "color": WHITE}),
       ("   ·   Head of Laboratory", {"color": PALE})]],
     size=13, color=PALE, font=BODY)
text(s, M, Inches(6.42), Inches(7.4), Inches(0.3),
     "lexuanbach@hcmut.edu.vn   ·   lexuanbach.github.io",
     size=11, color=PALE, font=BODY)

# ====================================================== 2 · MOTIVATION ===
s = slide(
    "This is the hook, and the reason the lab exists. Say it slowly.\n\n"
    "\"Two years ago AI suggested things. Now it ships them. It writes "
    "patches, it ranks security alerts, it drafts answers people act on. So "
    "ask the obvious question: how do you know it was right? Most of the "
    "time the honest answer is a benchmark score.\"\n\n"
    "Then land the turn: a score tells you what happened on someone else's "
    "test set last month. It is not a reason to trust tomorrow's answer. "
    "That gap is our research programme."
)
bg(s, WHITE)
rect(s, 0, 0, Pt(7), H, fill=CRIMSON)
text(s, M, Inches(1.35), Inches(10.6), Inches(0.4),
     [[("RESEARCH MOTIVATION", {})]], size=12, color=CRIMSON, bold=True)
text(s, M, Inches(1.95), Inches(11.0), Inches(2.6),
     [[("AI systems now write code, rank alerts\nand draft answers.",
        {"color": NAVY}),
       ("  Very little of that output\ncarries evidence that it is correct.",
        {"color": CRIMSON})]],
     size=32, font=HEAD, spacing=1.22)
rect(s, M, Inches(4.72), Inches(11.0), Pt(1), fill=RULE)
text(s, M, Inches(5.05), Inches(5.15), Inches(1.4),
     "A benchmark score summarises what happened last time. That is a "
     "useful thing to know. It is not a reason to trust the next answer.",
     size=14.5, color=MUTED, font=BODY, spacing=1.3)
text(s, M + Inches(5.85), Inches(5.05), Inches(5.15), Inches(1.4),
     [[("RAISE builds what is missing:", {"bold": True, "color": NAVY}),
       (" tests written to break a candidate patch, models that report "
        "their own error bars, and audits that say plainly where the "
        "evidence stops.", {"color": MUTED})]],
     size=14.5, font=BODY, spacing=1.3)
footer(s, 2)

# ====================================================== 3 · PRINCIPLES ===
s = slide(
    "Our three working principles. This is the lab's character, and it's "
    "what students actually learn here.\n\n"
    "On the first: I tell students that a model sounding sure of itself is "
    "not evidence of anything. Find an external check.\n"
    "On the second: before you report a number, know what it measures.\n"
    "On the third: 'inconclusive' is a real answer and we design for it. "
    "Reviewers respect it, and users can act on it."
)
bg(s, WHITE)
header(s, "Guiding principles", "Three commitments behind every project")
cards = [
    ("Evidence over confidence",
     "A model sounding certain proves nothing. We make every claim answer to "
     "something outside the model: a test that tries to break it, a solver, "
     "a statistical check."),
    ("Measurement that means something",
     "Before we report a number we ask what it measures and whether the data "
     "supports it. An evaluation is a measuring instrument. Instruments get "
     "calibrated."),
    ("Honest about limits",
     "We build systems that can answer “inconclusive”. Knowing where "
     "the evidence runs out is more useful, and more honest, than a "
     "confident guess."),
]
cw = (CW - Inches(0.5)) / 3
for i, (t, d) in enumerate(cards):
    x = M + i * (cw + Inches(0.25))
    rect(s, x, Inches(2.15), cw, Inches(3.5), fill=PANEL)
    rect(s, x, Inches(2.15), cw, Pt(3.5), fill=[CRIMSON, AMBER, NAVY2][i])
    text(s, x + Inches(0.4), Inches(2.55), Inches(0.8), Inches(0.5),
         f"0{i+1}", size=26, color=[CRIMSON, AMBER, NAVY2][i], bold=True,
         font=HEAD)
    text(s, x + Inches(0.4), Inches(3.20), cw - Inches(0.8), Inches(0.7), t,
         size=16.5, color=NAVY, bold=True, font=BODY, spacing=1.1)
    text(s, x + Inches(0.4), Inches(4.10), cw - Inches(0.8), Inches(1.4), d,
         size=12.5, color=MUTED, font=BODY, spacing=1.22)
footer(s, 3)

# ================================================== 4 · PROGRAMME MAP ===
s = slide(
    "The map of the lab. Give people the shape before the detail.\n\n"
    "\"We work across three pillars. Software engineering is the biggest, "
    "and that's agents, static analysis and reliability. AI and machine "
    "learning is the methods layer underneath it. Security is where both of "
    "those meet somebody who is actively trying to break them.\"\n\n"
    "Mention the fourth strand at the bottom if you have industry people in "
    "the room. It's our most applied work."
)
bg(s, NAVY)
header(s, "Research programme", "Three pillars, fourteen active directions",
       dark=True)
pillars = [
    ("Software\nEngineering", "5 directions",
     "Agents across the lifecycle · AI governance · Static analysis and "
     "testing · Agent reliability · Quantum program testing", CRIMSON),
    ("Artificial Intelligence\nand Machine Learning", "5 directions",
     "Evaluation and benchmarks · Uncertainty and selective prediction · "
     "Efficient AI systems · Adaptive vision · Retrieval and multilingual "
     "NLP", AMBER),
    ("Security", "4 directions",
     "Privacy in AI agents · Privacy-preserving recommendation · "
     "Vulnerability detection · Trustworthy AI-generated code",
     RGBColor(0x6C, 0x9B, 0xD1)),
]
cw = (CW - Inches(0.5)) / 3
for i, (name, cnt, body, col) in enumerate(pillars):
    x = M + i * (cw + Inches(0.25))
    rect(s, x, Inches(2.05), cw, Inches(3.75), fill=NAVY2)
    rect(s, x, Inches(2.05), cw, Pt(4), fill=col)
    text(s, x + Inches(0.38), Inches(2.45), cw - Inches(0.7), Inches(1.05),
         name, size=19, color=WHITE, font=HEAD, spacing=1.14)
    text(s, x + Inches(0.38), Inches(3.62), cw - Inches(0.7), Inches(0.28),
         cnt.upper(), size=10.5, color=col, bold=True, font=BODY)
    text(s, x + Inches(0.38), Inches(4.05), cw - Inches(0.7), Inches(1.6),
         body, size=12, color=PALE, font=BODY, spacing=1.3)
text(s, M, Inches(6.15), CW, Inches(0.5),
     [[("A fourth strand runs alongside these:  ",
        {"bold": True, "color": WHITE}),
       ("decision-making under uncertainty, where a forecast feeds a real "
        "optimisation step. ATM cash management and examination room "
        "allocation, both published this year.", {"color": PALE})]],
     size=12.5, font=BODY, spacing=1.25)
footer(s, 4, dark=True)

# ====================================================== 5 · PILLAR · SE ===
s = slide(
    "Pillar one, and our largest. Don't read all five. Pick two and go "
    "deep.\n\n"
    "AI governance is the easiest one for a non-specialist to grasp: a patch "
    "that passes the tests you already had has proved almost nothing, so we "
    "write new tests whose whole job is to break it.\n\n"
    "The quantum line usually surprises people. It's the oldest thread here, "
    "going back to POPL and ICALP, and it's why this lab reaches for proof "
    "before it reaches for a benchmark."
)
bg(s, WHITE)
header(s, "Pillar I", "Software Engineering")
rows(s, [
    ("AI Agents in Software Development",
     "Agents now touch requirements, code, tests and operations. We map what "
     "they actually do at each stage, then build evaluations that measure "
     "something real instead of a leaderboard delta."),
    ("AI Governance",
     "Someone should be able to ask why a system accepted AI-written code. "
     "We record what justified the decision. A patch that passes the existing "
     "tests has proved little, so we write new ones built to break it."),
    ("Static Analysis and Testing",
     "Nobody trusts an analyser that cries wolf. We rank alerts on evidence "
     "pulled from the repository and attach honest confidence to each, so "
     "the queue is worth working through."),
    ("Agent Reliability",
     "An agent is a stochastic system, so we model it as one. Fitting "
     "execution traces to Markov chains yields error bars and a "
     "goodness-of-fit test. A single benchmark score yields neither."),
    ("Testing for Quantum Programs",
     "You cannot debug a quantum program by running it and printing values. "
     "Correctness has to come from the type system and the logic, which is "
     "where our earlier work on separation logic pays off."),
], top=Inches(1.80), row_h=Inches(0.96), gap=Inches(0.045))
footer(s, 5)

# ==================================================== 6 · PILLAR · AIML ===
s = slide(
    "Pillar two, the methods layer. These are the techniques that make the "
    "software engineering work possible.\n\n"
    "Best one for a general audience is selective prediction: we build "
    "systems that decline to answer when they should, and we want that "
    "guarantee to hold on a finite sample, not on a threshold somebody "
    "tuned by hand.\n\n"
    "The efficiency question is a good one to put to industry people: five "
    "small models or one big one? You can't answer it until you fix the "
    "budget."
)
bg(s, WHITE)
header(s, "Pillar II", "Artificial Intelligence and Machine Learning")
rows(s, [
    ("AI Evaluation and Benchmarks",
     "Using a model as a judge turns it into a measuring instrument, and "
     "instruments need calibrating. We score judges against verifiable "
     "ground truth and watch for benchmarks going stale through "
     "contamination."),
    ("Uncertainty and Selective Prediction",
     "Some questions a system should decline. We build that refusal on "
     "finite-sample guarantees, set thresholds per group so fairness is not "
     "quietly traded away, and let error costs decide how cautious to be."),
    ("Efficient AI Systems",
     "Five small models or one large one? Nobody can answer that without "
     "fixing the budget first. We run the comparison at equal cost, and ask "
     "how far aggregation goes before voting theory bites."),
    ("Adaptive Computer Vision",
     "A frozen encoder meets a satellite sensor nobody trained it on. We "
     "adapt it at test time and without labels, reweighting spectral bands "
     "and routing adapters to the layers that need them."),
    ("Retrieval and Multilingual NLP",
     "Legal text is full of amendments, citations and definitions that flat "
     "retrieval flattens away. We walk the knowledge graph alongside "
     "semantic search so that structure survives into the answer."),
], top=Inches(1.80), row_h=Inches(0.96), gap=Inches(0.045))
footer(s, 6)

# ================================================ 7 · PILLAR · SECURITY ===
s = slide(
    "Pillar three. This is where our work meets an actual adversary, and it "
    "lands hardest with industry.\n\n"
    "Lead with privacy in AI agents. It's concrete and slightly alarming: "
    "persistent memory is the least audited part of an agent deployment, and "
    "what one user's session writes into a shared memory, another user's "
    "session can read back. We built a benchmark so this gets measured "
    "instead of argued about.\n\n"
    "Note the philosophy in vulnerability detection: we tune for a low "
    "false-discovery rate, not maximal recall. A ranked queue should mean "
    "what it says."
)
bg(s, WHITE)
header(s, "Pillar III", "Security")
rows(s, [
    ("Privacy in AI Agents",
     "Persistent memory is the least audited part of an agent deployment. "
     "What one user's session writes, another user's session can read back. "
     "We built a benchmark and an audit suite so this gets measured, not "
     "argued about."),
    ("Privacy-Preserving Recommendation",
     "A deployed recommender is asked for differential privacy, for "
     "calibration inside every group, and for the right to stay silent. "
     "Optimise these one at a time and each eats the others. We treat them "
     "as one constrained problem."),
    ("Vulnerability Detection",
     "Every hypothesis is checked against what the repository actually "
     "shows, then marked Verified, Rejected or Inconclusive. We would rather "
     "report ten real bugs than a hundred maybes, so we tune for false "
     "discovery, not recall."),
    ("Trustworthy AI-Generated Code",
     "Translated and repaired code has no author to ask. We put it through "
     "static analysis and generated tests, so accepting it rests on "
     "something firmer than the model's own say-so."),
], top=Inches(1.88), row_h=Inches(1.05), gap=Inches(0.10))
footer(s, 7)

# ================================================ 8 · REPRESENTATIVE WORK ===
s = slide(
    "Four concrete projects. This is the slide that makes the principles "
    "real, so pick whichever fits the room.\n\n"
    "Software engineering people: TraceToChain. Security people: "
    "GraphMemShield or GraphLedger. Data and ML people: LTM-RAG.\n\n"
    "The LTM-RAG figure is the most quotable number in the deck. Attack "
    "success falls from 94.8% to 14.0%."
)
bg(s, WHITE)
header(s, "Representative work", "Four selected projects")
projects = [
    ("TraceToChain", "ISSRE 2026",
     "Turns agent execution traces into an absorbing Markov chain, then "
     "refuses to report any number the statistics do not support. Tested on "
     "2,459 real agent episodes.", CRIMSON),
    ("GraphLedger", "ASE 2026",
     "Finds vulnerabilities across a whole repository and makes every "
     "finding show its work against real dependency and taint structure.",
     AMBER),
    ("GraphMemShield", "CIKM 2026",
     "Measures what leaks between users in an agent's shared memory. "
     "Privacy stops being an argument and becomes a number.",
     RGBColor(0x6C, 0x9B, 0xD1)),
    ("LTM-RAG", "ICDM 2026",
     "Treats poisoning of retrieval-augmented generation as something that "
     "happens over time, not one query at a time. Attack success falls from "
     "94.8% to 14.0%.", RGBColor(0x7E, 0xB0, 0x9B)),
]
cw2 = (CW - Inches(0.3)) / 2
for i, (name, venue, desc, col) in enumerate(projects):
    x = M + (i % 2) * (cw2 + Inches(0.3))
    y = Inches(2.02) + (i // 2) * Inches(2.30)
    rect(s, x, y, cw2, Inches(2.05), fill=PANEL)
    rect(s, x, y, Pt(3.5), Inches(2.05), fill=col)
    text(s, x + Inches(0.34), y + Inches(0.26), cw2 - Inches(0.7), Inches(0.4),
         name, size=20, color=NAVY, font=HEAD)
    text(s, x + Inches(0.34), y + Inches(0.72), cw2 - Inches(0.7), Inches(0.26),
         venue.upper(), size=10, color=col, bold=True, font=BODY)
    text(s, x + Inches(0.34), y + Inches(1.06), cw2 - Inches(0.7), Inches(0.9),
         desc, size=12, color=MUTED, font=BODY, spacing=1.22)
footer(s, 8)

# ================================================== 9 · CURRENT DIRECTIONS ===
s = slide(
    "This is the 'what's next' slide, and it's the one that starts "
    "conversations. Deliberately no project names here: a lot of it is under "
    "review.\n\n"
    "\"The first wave of agent research asked what these systems can do. "
    "That question is mostly settled. What interests us now is second-order: "
    "not whether an agent can act, but when it should, and on whose "
    "evidence.\"\n\n"
    "Walk the five questions. Each has two or three papers behind it, either "
    "in preparation or under submission. If someone asks for detail on one, "
    "that's the collaboration conversation worth having after the talk.\n\n"
    "The independence point in question two is the one people remember: ten "
    "sources that copied each other are not ten pieces of evidence."
)
bg(s, WHITE)
header(s, "Current directions", "Ongoing research")
text(s, M, Inches(1.80), Inches(11.0), Inches(0.62),
     [[("The first wave of agent research asked what these systems can do. "
        "The questions we find interesting now are second-order: ",
        {"color": MUTED}),
       ("not whether an agent can act, but when it should, and on whose "
        "evidence.", {"color": NAVY, "bold": True})]],
     size=14, font=BODY, spacing=1.3)

qs = [
    ("When should\nit search?",
     "Retrieval treated as a decision with a cost. Look something up only "
     "while the next lookup is worth more than it costs, and ask the user "
     "when asking beats guessing.", CRIMSON),
    ("Whom should\nit trust?",
     "Ten sources that copied one another are not ten pieces of evidence. We "
     "weigh evidence by independence, and stop untrusted content from ever "
     "carrying authority.", AMBER),
    ("What should\nit remember?",
     "Facts expire. We work on memory that tracks how long something stays "
     "true, and notices when fresh evidence contradicts what it stored.",
     RGBColor(0x6C, 0x9B, 0xD1)),
    ("How does it\nsurvive change?",
     "Interfaces move, distributions shift, benchmarks go stale. A system "
     "should notice when it is running outside the conditions it was "
     "validated on.", RGBColor(0x7E, 0xB0, 0x9B)),
    ("Can it show\nits work?",
     "Provenance you could hand to an auditor, tracing an output back "
     "through the data and the decisions that produced it. Sometimes with "
     "cryptographic guarantees.", RGBColor(0x9B, 0x8E, 0xC4)),
]
n = len(qs)
gapx = Inches(0.14)
cwq = (CW - gapx * (n - 1)) / n
for i, (q, d, col) in enumerate(qs):
    x = M + i * (cwq + gapx)
    rect(s, x, Inches(2.72), cwq, Inches(3.30), fill=PANEL)
    rect(s, x, Inches(2.72), cwq, Pt(3.5), fill=col)
    text(s, x + Inches(0.24), Inches(3.00), cwq - Inches(0.48), Inches(0.78),
         q, size=15, color=NAVY, font=HEAD, spacing=1.14)
    text(s, x + Inches(0.24), Inches(3.92), cwq - Inches(0.48), Inches(1.9),
         d, size=10.5, color=MUTED, font=BODY, spacing=1.22)
text(s, M, Inches(6.22), CW, Inches(0.4),
     "Each question has work in preparation or under submission across "
     "search, security, data mining and machine learning venues.",
     size=11.5, color=MUTED, font=BODY)
footer(s, 9)

# ================================================== 9 · PUBLICATION RECORD ===
s = slide(
    "Credibility slide. State it plainly and move on.\n\n"
    "\"Twenty-eight publications, nineteen of them this year. The recent "
    "work sits at A and A-star venues in software engineering and data "
    "mining. The earlier work is the formal methods foundation: POPL, "
    "ICALP, ESOP, NeurIPS.\"\n\n"
    "The point worth making: the lab is new, but the methods are not. The "
    "rigour comes out of a decade in program verification."
)
bg(s, NAVY)
header(s, "Publication record", "A new laboratory built on established "
        "methods", dark=True)
stats = [("28", "Publications"), ("19", "Papers in 2026"),
         ("14", "Research directions")]
sw = (CW - Inches(0.15) * (len(stats) - 1)) / len(stats)
for i, (v, l) in enumerate(stats):
    x = M + i * (sw + Inches(0.15))
    rect(s, x, Inches(1.98), sw, Inches(1.42), fill=NAVY2)
    text(s, x, Inches(2.22), sw, Inches(0.7), v, size=44, color=WHITE,
         font=HEAD, align=PP_ALIGN.CENTER)
    text(s, x, Inches(2.95), sw, Inches(0.3), l.upper(), size=10,
         color=CRIMSON, bold=True, font=BODY, align=PP_ALIGN.CENTER)

text(s, M, Inches(3.75), CW, Inches(0.3),
     [[("RECENT — SOFTWARE ENGINEERING, DATA MINING, NLP, SECURITY", {})]],
     size=10.5, color=CRIMSON, bold=True, font=BODY)
text(s, M, Inches(4.12), CW, Inches(0.5),
     "ICDM 2026  ×2   ·   EMNLP 2026   ·   CIKM 2026   ·   ASE 2026  ×2   ·   "
     "ISSRE 2026   ·   ACIVS 2026  ×3   ·   ACIIDS 2026  ×2   ·   IEA/AIE 2026  ×4",
     size=14.5, color=WHITE, font=BODY, spacing=1.3)
rect(s, M, Inches(4.92), CW, Pt(1), fill=RGBColor(0x3A, 0x3F, 0x55))
text(s, M, Inches(5.20), CW, Inches(0.3),
     [[("FOUNDATIONS — PROGRAM VERIFICATION, LOGIC, COMPLEXITY", {})]],
     size=10.5, color=AMBER, bold=True, font=BODY)
text(s, M, Inches(5.57), CW, Inches(0.5),
     "NeurIPS 2024   ·   POPL 2022   ·   ICECCS 2020   ·   ICALP 2019   ·   "
     "APLAS 2018   ·   ESOP 2018   ·   ICFEM 2017   ·   FSTTCS 2016",
     size=14.5, color=WHITE, font=BODY, spacing=1.3)
footer(s, 10, dark=True)

# ==================================================== 10 · LABORATORY TEAM ===
s = slide(
    "Who we are. Keep the biography brief. The students are the interesting "
    "part.\n\n"
    "\"My PhD was at NUS with Aquinas Hobor, with mentorship from Anthony "
    "Lin, then a postdoc at NTU Singapore with Luke Ong. That background is "
    "program verification and logic, which is why this lab reaches for proof "
    "and measurement before it reaches for a benchmark.\"\n\n"
    "The names shown are a selection of students with published work. Say "
    "that out loud, because the lab supervises others who aren't listed "
    "here. The promise worth making: everyone who joins gets a problem of "
    "their own and publishes it."
)
bg(s, WHITE)
header(s, "Laboratory team", "Supervision and members")
rect(s, M, Inches(1.98), Inches(5.0), Inches(4.05), fill=NAVY)
rect(s, M, Inches(1.98), Inches(5.0), Pt(4), fill=CRIMSON)
text(s, M + Inches(0.42), Inches(2.35), Inches(4.2), Inches(0.5),
     "Dr. Xuan-Bach Le", size=25, color=WHITE, font=HEAD)
text(s, M + Inches(0.42), Inches(2.92), Inches(4.2), Inches(0.3),
     "HEAD OF LABORATORY · LECTURER, HCMUT", size=10.5, color=CRIMSON,
     bold=True, font=BODY)
text(s, M + Inches(0.42), Inches(3.42), Inches(4.2), Inches(2.2),
     [["PhD, National University of Singapore, in program verification, "
       "under Aquinas Hobor and with mentorship from Anthony W. Lin."],
      ["Postdoctoral researcher, NTU Singapore, with Prof. Luke Ong."],
      ["Double degree in Computer Science and Pure Mathematics, NUS."]],
     size=12.5, color=PALE, font=BODY, spacing=1.25, space_after=9)

sx = M + Inches(5.3)
sw2 = CW - Inches(5.3)
text(s, sx, Inches(2.02), sw2, Inches(0.3),
     [[("CURRENT STUDENTS · SELECTED", {})]],
     size=10.5, color=CRIMSON, bold=True, font=BODY)
cur = [
    ("Thanh-Hai Tran", "Visual anomaly detection and risk-aware recovery"),
    ("Thi-Hong-Cuc Le", "Repository-level vulnerability detection and explainable AI"),
    ("Hoang-Quoc-Bao Hua", "LLM agents for static analysis and assumption validation"),
    ("Huu-Vu-Phuong Tran", "Poison-resilient retrieval-augmented generation"),
    ("Duc-Thuan Mai", "Longitudinal trust mining for retrieval pipelines"),
    ("Van-Truong-Thinh Nguyen", "Depth-aware vision transformers for air quality"),
    ("Tien-Anh Nguyen", "Test-time adaptation of vision-language models"),
    ("Minh-Hieu Le", "Examination room optimisation, greedy and exact methods"),
]
y = Inches(2.34)
for n, t in cur:
    text(s, sx, y, sw2, Inches(0.22),
         [[(n, {"bold": True, "color": NAVY, "size": 12}),
           ("   " + t, {"color": MUTED, "size": 10.5})]], font=BODY,
         spacing=1.05)
    y += Inches(0.35)
rect(s, sx, y + Inches(0.04), sw2, Pt(1), fill=RULE)
text(s, sx, y + Inches(0.22), sw2, Inches(0.3),
     [[("RECENT GRADUATES · SELECTED", {})]], size=10.5, color=AMBER,
     bold=True, font=BODY)
text(s, sx, y + Inches(0.56), sw2, Inches(0.6),
     "Khoa Phan  ·  Phong Chung  ·  Huu-Thanh Phan  ·  Huynh Duy Khang",
     size=11.5, color=MUTED, font=BODY, spacing=1.2)
text(s, sx, y + Inches(0.92), sw2, Inches(0.3),
     [[("COLLABORATORS", {})]], size=10.5, color=RGBColor(0x6C, 0x9B, 0xD1),
     bold=True, font=BODY)
text(s, sx, y + Inches(1.26), sw2, Inches(0.4),
     "Phat T. Tran-Truong (HCMUT)  ·  Son Ha Xuan (RMIT)",
     size=11.5, color=MUTED, font=BODY, spacing=1.2)
footer(s, 11)

# ====================================================== 11 · COLLABORATION ===
s = slide(
    "The close. Make the ask explicit and give people one thing to do.\n\n"
    "For students: 'I'm looking for people who want to do careful work. "
    "Email me and tell me what you find interesting.'\n\n"
    "For collaborators, name the two things we most want: real agent traces "
    "and real alert queues. Our methods are only worth as much as the data "
    "they meet, and that data sits inside companies, not universities.\n\n"
    "Leave this slide up during questions."
)
bg(s, NAVY)
rect(s, 0, 0, Pt(7), H, fill=CRIMSON)
text(s, M, Inches(0.95), CW, Inches(0.35),
     [[("COLLABORATION AND ADMISSIONS", {})]], size=12, color=CRIMSON,
     bold=True, font=BODY)
text(s, M, Inches(1.42), Inches(10.4), Inches(1.0),
     "Careful work, honest about its limits,\nand worth publishing.",
     size=34, color=WHITE, font=HEAD, spacing=1.2)

boxes = [
    ("Prospective students",
     "Master's and undergraduate projects across all three pillars. You will "
     "own a problem, publish it, and learn to defend a claim with evidence.",
     CRIMSON),
    ("Academic collaborators",
     "We are looking for real agent traces and real static-analysis alert "
     "queues. Our methods are only worth as much as the data they meet.",
     AMBER),
    ("Industry partners",
     "Reliability audits, leakage assessments and evaluation design for "
     "deployed AI systems. Measurement you can put in front of a regulator.",
     RGBColor(0x6C, 0x9B, 0xD1)),
]
cw3 = (CW - Inches(0.5)) / 3
for i, (t, d, col) in enumerate(boxes):
    x = M + i * (cw3 + Inches(0.25))
    rect(s, x, Inches(3.20), cw3, Inches(2.05), fill=NAVY2)
    rect(s, x, Inches(3.20), cw3, Pt(3.5), fill=col)
    text(s, x + Inches(0.32), Inches(3.55), cw3 - Inches(0.65), Inches(0.35),
         t, size=16, color=WHITE, bold=True, font=BODY)
    text(s, x + Inches(0.32), Inches(4.02), cw3 - Inches(0.65), Inches(1.2),
         d, size=12, color=PALE, font=BODY, spacing=1.25)

rect(s, M, Inches(5.72), CW, Pt(1), fill=RGBColor(0x3A, 0x3F, 0x55))
text(s, M, Inches(6.05), CW, Inches(0.8),
     [[("Dr. Xuan-Bach Le", {"bold": True, "color": WHITE, "size": 15}),
       ("    lexuanbach@hcmut.edu.vn", {"color": PALE, "size": 13}),
       ("    ·    lexuanbach.github.io", {"color": PALE, "size": 13})]],
     font=BODY)
text(s, M, Inches(6.45), CW, Inches(0.4),
     "Faculty of Computer Science and Engineering, HCMUT · "
     "268 Ly Thuong Kiet, District 10, Ho Chi Minh City",
     size=11, color=PALE, font=BODY)


# --------------------------------------------------------------- narrative ---
# Full spoken script, one entry per slide, in delivery order.
NARRATIVE = [

# 1 · Title
"""Thanks for having me. I'm Xuan-Bach Le. I teach at Ho Chi Minh City
University of Technology, and I run a small research lab there called RAISE,
which stands for Reasoning in Artificial Intelligence and Software Engineering.

That is a long name for a simple idea. AI systems now produce a great deal of
software. They write patches, they rank security alerts, they draft answers
that people act on. Almost none of that output arrives with any evidence that
it is correct. We build the layer that checks it.

I will keep this to about fifteen minutes. I want to tell you why the problem
matters, walk you through the three pillars we work across, show you four
projects we published this year, and then spend a little time on what we are
doing right now, because that is usually where a collaboration starts.

[Delivery: warm and unhurried. Do not rush the opening; this is the only
moment people decide whether to listen.]""",

# 2 · Research motivation
"""Let me start with the problem, because everything else follows from it.

Two years ago AI suggested things to us. Now it ships them. It writes the patch
that goes into your repository. It ranks the security alerts your team works
through on Monday morning. It drafts the answer somebody acts on without ever
opening the source.

So ask the obvious question. How do you know it was right?

Most of the time the honest answer is a benchmark score. Something like: this
model scores seventy-two percent on a public test set. That is a genuinely
useful thing to know. But look at what it actually tells you. It summarises
what happened, on somebody else's data, at some point in the past. It is not a
reason to trust the next answer, and the next answer is the only one you
actually care about.

That gap is the entire research programme. We build the things missing from
that picture. Tests written specifically to break a candidate patch. Models
that report their own error bars instead of a single number. Audits that say
plainly where the evidence stops and the guessing begins.

[Delivery: slow down on "how do you know it was right?" and let it sit for a
beat before answering.]""",

# 3 · Guiding principles
"""Before the research itself, three commitments. These shape how we work, and
honestly they are what students actually learn in the lab.

The first is evidence over confidence. A model sounding certain proves nothing
at all. Language models are fluent by construction, and fluency reads to us as
confidence. So we make every claim answer to something outside the model. A
test that tries to break it. A solver. A statistical check. Something with the
power to say no.

The second is measurement that means something. Before we report a number we
ask what it measures and whether the data supports it. An evaluation is a
measuring instrument, and instruments get calibrated. A surprising amount of
our work is just taking that seriously.

The third is being honest about limits. We build systems that can answer
"inconclusive". That sounds like a weakness and it is the opposite. Knowing
where your evidence runs out is far more useful to whoever depends on you than
a confident guess. Reviewers respect it. Users can act on it.

If you take one thing from this talk, take the third one.""",

# 4 · Research programme
"""Here is the shape of the lab before I go into any detail.

Three pillars. Software engineering is the largest, covering agents across the
development lifecycle, static analysis, and reliability. Artificial
intelligence and machine learning is the methods layer underneath it, where the
evaluation, uncertainty and vision work lives. Security is where both of those
meet somebody who is actively trying to break them.

Fourteen active directions across the three. I will not walk through all
fourteen. What I would rather you take away is that these are not three
separate labs sharing a corridor. The same question runs through all of them,
which is how do you know this thing is right, and the pillars are simply three
places where that question gets asked.

There is a fourth strand at the bottom that does not fit the pattern as
neatly. That is decision-making under uncertainty, where a forecast feeds a
real optimisation step. Two examples published this year: cash management for
ATM networks, and examination room allocation at our own university. It is our
most applied work, and it is often the fastest way into a conversation with
industry.""",

# 5 · Pillar I — Software Engineering
"""Pillar one, and our largest. I will not read all five, so let me take two.

AI governance is the easiest to grasp without a software background. When a
system accepts AI-written code, someone should be able to ask why. Not "the
model was confident", but an actual record of what justified the decision. And
here is the part people find surprising: a patch that passes the tests you
already had has proved almost nothing, because those tests were written before
the patch existed. So we generate new ones whose entire job is to break the
candidate, and we let the result decide instead of the model.

Agent reliability is the second. An agent is a stochastic system. It retries,
it wanders, sometimes it succeeds on the fourth attempt. So we model it as one.
Fitting execution traces to a Markov chain gives you error bars and a
goodness-of-fit test. A single benchmark score gives you neither.

The quantum line at the bottom usually surprises people. It is the oldest
thread here, going back to POPL and ICALP. You cannot debug a quantum program
by running it and printing values, so correctness has to come from the type
system and the logic. That is the tradition this lab reasons from.""",

# 6 · Pillar II — AI and ML
"""Pillar two is the methods layer. These are the techniques that make the
software engineering work possible.

The one I would point a general audience to is uncertainty and selective
prediction. The idea is simple: some questions a system should decline to
answer. Building that well is harder than it sounds. You want the refusal to
rest on a guarantee that holds on a finite sample, not on a threshold somebody
tuned by hand until the demo looked good. We also set those thresholds per
group, because a system that becomes more cautious for only some populations
has quietly traded away fairness.

The efficiency line is a good question to put to anyone running AI in
production. Five small models or one large one? People argue about this
constantly and almost nobody fixes the budget first. Once you compare at equal
cost the answer changes, and it changes by task.

The evaluation work at the top is the one I would flag for researchers. The
moment you use a model as a judge, you have turned it into a measuring
instrument, and instruments need calibrating against something you can
verify.""",

# 7 · Pillar III — Security
"""Pillar three is where our work meets an actual adversary. This is the
pillar that lands hardest with industry.

Start with privacy in AI agents, because it is concrete and slightly alarming.
Persistent memory is the least audited part of any agent deployment. Everyone
scrutinises the model and the prompt. Almost nobody audits the memory. And what
one user's session writes into a shared memory, another user's session can read
back. We built a benchmark and an audit suite so that this gets measured
instead of argued about in a meeting.

Vulnerability detection is where I would point your security team. Every
hypothesis is checked against what the repository actually shows, and then
marked Verified, Rejected, or Inconclusive. Notice the philosophy. We would
rather report ten real bugs than a hundred maybes. We tune for a low
false-discovery rate, not maximal recall, because a ranked queue should mean
what it says. Once your engineers stop trusting the queue, the tool is
worthless no matter how good its recall number looks.""",

# 8 · Representative work
"""Four projects from this year, to make all of that concrete.

TraceToChain takes agent execution traces and turns them into an absorbing
Markov chain. The interesting part is not the model, it is that the system
refuses to report any number the statistics do not support. We tested it on
2,459 real agent episodes drawn from SWE-bench and tau-bench.

GraphLedger finds vulnerabilities across a whole repository and makes every
finding show its work against real dependency and taint structure.

GraphMemShield is the memory leakage work I just mentioned. It measures what
leaks between users in an agent's shared memory, so privacy stops being an
argument and becomes a number.

LTM-RAG treats poisoning of retrieval systems as something that happens over
time rather than one query at a time, which is how these attacks actually work.
That reframing is worth a number: answer-level attack success falls from
ninety-five percent to fourteen.

Pick whichever of these is closest to what you do and I am happy to go deeper
afterwards.""",

# 9 · Current directions
"""This is the part I am most interested in talking about. There are no
project names on this slide on purpose, because most of it is under review.

Here is the framing. The first wave of agent research asked what these systems
can do. That question is largely settled; they can do a great deal. What
interests us now is second-order. Not whether an agent can act, but when it
should, and on whose evidence.

Five questions. When should it search? Treat retrieval as a decision with a
cost, and look something up only while the next lookup is worth more than it
costs. Sometimes the right move is to ask the user instead of guessing.

Whom should it trust? This is my favourite. Ten sources that copied one another
are not ten pieces of evidence. So we weigh evidence by independence, and we
stop untrusted content from ever carrying authority.

What should it remember? Facts expire, and memory ought to know that.

How does it survive change? Interfaces move, distributions shift, benchmarks go
stale. A system should notice when it is running outside the conditions it was
validated on.

And can it show its work? Provenance you could hand to an auditor.""",

# 10 · Publication record
"""Briefly, so you know the lab delivers.

Twenty-eight publications, nineteen of them this year. The recent work sits at
A and A-star venues in software engineering and data mining: ICDM, CIKM, ASE,
ISSRE.

The bottom row is the part I would draw your attention to. POPL, ICALP, ESOP,
FSTTCS, NeurIPS. That is program verification, logic and complexity, and it is
older work.

The point is this. The lab is new. The methods are not. Everything I have shown
you today rests on roughly a decade of program verification, and that is why we
reach for a proof or a statistical certificate before we reach for a
leaderboard. It is an unusual foundation for a group working on LLM agents, and
I think it is our main advantage.""",

# 11 · Laboratory team
"""A word about who actually does this work.

My own background is program verification. PhD at the National University of
Singapore under Aquinas Hobor, with mentorship from Anthony Lin, then a
postdoc at NTU Singapore with Luke Ong. That is where the instinct for proof
and measurement comes from.

The names on the right are a selection of students with published work, and I
want to be clear that it is a selection. The lab supervises others who are not
on this slide.

What I would rather you notice is the range. Vulnerability detection, agent
reliability, retrieval security, computer vision, optimisation. These are
master's and undergraduate students, and they are first authors.

That is deliberate, and it is the promise I make when somebody joins. You get a
problem of your own, you run it, and you publish it. I do not hand out
subtasks. It is a slower way to run a lab and I believe it produces better
researchers.""",

# 12 · Collaboration and admissions
"""Three ways to work with us, and then I will take questions.

If you are a student, or you teach students who might fit: I am looking for
people who want to do careful work. The line at the top of this slide is the
actual standard. Careful, honest about its limits, and worth publishing.
Projects are open across all three pillars. Email me and tell me what you find
interesting.

If you are an academic collaborator, here is my concrete ask. We need real
agent traces and real static-analysis alert queues. Our methods are only worth
as much as the data they meet, and that data mostly sits inside companies
rather than universities. If you have it, or you know who does, that is the
most useful thing you could offer us.

And for industry: reliability audits, leakage assessments, evaluation design.
Measurement you could put in front of a regulator.

My email is on the slide. Thank you, and I am happy to take questions.

[Delivery: leave this slide up for the whole Q&A.]""",
]

for _slide, _script in zip(prs.slides, NARRATIVE):
    _tf = _slide.notes_slide.notes_text_frame
    _tf.text = "\n\n".join(
        " ".join(_para.split())
        for _para in _script.strip().split("\n\n")
    )

import sys
out = sys.argv[1]
prs.save(out)
print(f"saved {out} · {len(prs.slides._sldIdLst)} slides")
